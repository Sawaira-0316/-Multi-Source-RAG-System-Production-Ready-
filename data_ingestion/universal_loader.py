# data_ingestion/universal_loader.py

import os
from pathlib import Path
from typing import List

import pandas as pd
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader
from langchain_core.documents import Document


def _load_txt(path: Path) -> List[Document]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    return [Document(page_content=text, metadata={"source": str(path), "type": "txt"})]


def _load_md(path: Path) -> List[Document]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    return [Document(page_content=text, metadata={"source": str(path), "type": "md"})]


def _load_pdf(path: Path) -> List[Document]:
    reader = PdfReader(str(path))
    docs: List[Document] = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        docs.append(
            Document(
                page_content=text,
                metadata={"source": str(path), "page": i, "type": "pdf"},
            )
        )
    return docs


def _load_docx(path: Path) -> List[Document]:
    doc = DocxDocument(str(path))
    text = "\n".join(p.text for p in doc.paragraphs)
    return [Document(page_content=text, metadata={"source": str(path), "type": "docx"})]


def _load_csv(path: Path) -> List[Document]:
    df = pd.read_csv(path)
    text = df.to_markdown(index=False)
    return [Document(page_content=text, metadata={"source": str(path), "type": "csv"})]


def _load_excel(path: Path) -> List[Document]:
    sheets = pd.read_excel(path, sheet_name=None)
    docs: List[Document] = []
    for name, df in sheets.items():
        text = f"# Sheet: {name}\n" + df.to_markdown(index=False)
        docs.append(
            Document(
                page_content=text,
                metadata={"source": str(path), "sheet": name, "type": "excel"},
            )
        )
    return docs


def _load_pptx(path: Path) -> List[Document]:
    prs = Presentation(str(path))
    texts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    full_text = "\n".join(texts)
    return [Document(page_content=full_text, metadata={"source": str(path), "type": "pptx"})]


def _load_html(path: Path) -> List[Document]:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(raw, "html.parser")
    text = soup.get_text(separator="\n")
    return [Document(page_content=text, metadata={"source": str(path), "type": "html"})]


EXTENSION_LOADERS = {
    ".txt": _load_txt,
    ".md": _load_md,
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".csv": _load_csv,
    ".xlsx": _load_excel,
    ".xls": _load_excel,
    ".pptx": _load_pptx,
    ".html": _load_html,
    ".htm": _load_html,
}


def load_file(path: Path) -> List[Document]:
    ext = path.suffix.lower()
    loader = EXTENSION_LOADERS.get(ext)
    if loader is None:
        # unsupported – skip silently
        return []
    return loader(path)


def load_corpus_from_paths(paths: List[str]) -> List[Document]:
    """
    Walk through all given paths (files/folders) and load supported formats
    into a list of LangChain Documents.
    """
    docs: List[Document] = []

    for p in paths:
        base = Path(p).expanduser()
        if not base.exists():
            continue

        if base.is_file():
            docs.extend(load_file(base))
        else:
            for root, _, files in os.walk(base):
                for fname in files:
                    fpath = Path(root) / fname
                    docs.extend(load_file(fpath))

    return docs
